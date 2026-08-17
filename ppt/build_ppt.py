# -*- coding: utf-8 -*-
"""
ppt/build_ppt.py —— 生成《从特斯拉到时空工程》汇报 PPTX（nature-paper2ppt 工作流）
论文类型：review（evidence-map 弧线）
输出：ppt/从特斯拉到时空工程-汇报.pptx + ppt/qa_report.md（QA 在运行审计后写入）
"""
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

try:
    sys.stdout.reconfigure(encoding="utf-8")
except Exception:
    pass

ROOT = Path(__file__).parent.parent
FIG = ROOT / "figures"
OUT = ROOT / "ppt"
OUT.mkdir(exist_ok=True)

# ---------- 设计常量 ----------
W, H = Inches(13.333), Inches(7.5)
DARK = RGBColor(0x1B, 0x1B, 0x2B)
ACCENT = RGBColor(0xC0, 0x39, 0x2B)
BLUE = RGBColor(0x29, 0x80, 0xB9)
GREEN = RGBColor(0x27, 0xAE, 0x60)
PURPLE = RGBColor(0x8E, 0x44, 0xAD)
GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT = RGBColor(0xF5, 0xF5, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "Microsoft YaHei"

prs = Presentation()
prs.slide_width = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]


def set_run(run, text, size=16, bold=False, color=DARK, font=FONT):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    rPr = run._r.get_or_add_rPr()
    from pptx.oxml.ns import qn
    from lxml import etree
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = etree.SubElement(rPr, qn("a:ea"))
    ea.set("typeface", font)


def add_text(slide, x, y, w, h, lines, align=PP_ALIGN.LEFT):
    """lines: list of (text, size, bold, color) 或 (text, size) 简写"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        text, size, bold, color = (line + (DARK,) if len(line) == 3 else line)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        set_run(p.add_run(), text, size, bold, color)
    return tb


def add_title(slide, text, num=None, color=ACCENT):
    """统一的标题条"""
    add_text(slide, Inches(0.55), Inches(0.28), Inches(12.2), Inches(0.7),
             [(text, 25, True, color)])
    # 标题下划线
    ln = slide.shapes.add_shape(1, Inches(0.58), Inches(0.98), Inches(2.4), Pt(2.2))
    ln.fill.solid()
    ln.fill.fore_color.rgb = color
    ln.line.fill.background()
    if num is not None:
        add_text(slide, Inches(12.35), Inches(0.15), Inches(0.8), Inches(0.5),
                 [(num, 13, True, GRAY)], PP_ALIGN.RIGHT)


def add_footer(slide, src="", page=0):
    add_text(slide, Inches(0.55), Inches(7.05), Inches(10.5), Inches(0.35),
             [(src, 9, False, GRAY)])
    add_text(slide, Inches(12.3), Inches(7.05), Inches(0.8), Inches(0.35),
             [(str(page), 9, False, GRAY)], PP_ALIGN.RIGHT)


def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def add_pic(slide, name, x, y, w=None, h=None):
    return slide.shapes.add_picture(str(FIG / name), x, y, width=w, height=h)


def add_table(slide, rows_data, x, y, w, h, col_widths=None, header=True, font_size=12):
    n_rows, n_cols = len(rows_data), len(rows_data[0])
    shape = slide.shapes.add_table(n_rows, n_cols, x, y, w, h)
    table = shape.table
    if col_widths:
        for j, cw in enumerate(col_widths):
            table.columns[j].width = cw
    for i, row in enumerate(rows_data):
        for j, cell_text in enumerate(row):
            cell = table.cell(i, j)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            bold = header and i == 0
            color = WHITE if (header and i == 0) else DARK
            set_run(p.add_run(), str(cell_text), font_size, bold, color)
            if header and i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = BLUE
            elif i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT
    return shape


def new_slide():
    return prs.slides.add_slide(BLANK)


def dark_bg(slide):
    rect = slide.shapes.add_shape(1, 0, 0, W, H)
    rect.fill.solid()
    rect.fill.fore_color.rgb = DARK
    rect.line.fill.background()
    return rect


# ================= Slide 1: 封面 =================
s = new_slide()
dark_bg(s)
add_text(s, Inches(1.2), Inches(2.2), Inches(11), Inches(1.4),
         [("从特斯拉到时空工程", 44, True, WHITE)], PP_ALIGN.CENTER)
add_text(s, Inches(1.2), Inches(3.45), Inches(11), Inches(0.8),
         [("统一场、量子现实与可编程时空 · 论文汇报", 22, False, RGBColor(0xCC, 0xCC, 0xDD))],
         PP_ALIGN.CENTER)
add_text(s, Inches(1.2), Inches(4.6), Inches(11), Inches(0.6),
         [("作者 MuningAn · 机构 PlanetarySystem · 2026", 15, False, RGBColor(0x99, 0x99, 0xAA))],
         PP_ALIGN.CENTER)
add_text(s, Inches(1.2), Inches(5.4), Inches(11), Inches(0.6),
         [("90 节全稿 + 19 行判决矩阵 + 16 份可复现脚本 · github.com/everest-an/Antigravity",
           13, False, RGBColor(0x88, 0x88, 0x99))], PP_ALIGN.CENTER)
add_notes(s, "开场：一句话定位。这不是一篇'特斯拉有没有反重力'的考证文，而是一张把'统一/反重力'所有说法变成可检验命题的地图。今天汇报按证据地图展开：为什么现在重要、三个主题、争议、判决矩阵、未来方向。")

# ================= Slide 2: 为什么现在重要 =================
s = new_slide()
add_title(s, "为什么这个主题现在重要", "02")
add_text(s, Inches(0.55), Inches(1.25), Inches(6.6), Inches(5.4), [
    ("桌面量子引力正处于爆发期", 20, True, BLUE),
    ("2017 年 BMV 提出'引力中介纠缠'实验后，这一领域在 2024–2026 年同时成熟了三条硬件路线：", 14, False, DARK),
    ("· 大质量叠加：纳米金刚石 Stern-Gerlach 干涉仪（Folman 组七份技术笔记，40 kHz 阱、UV 中和全部到位）", 13, False, DARK),
    ("· 时空涨落探测：GQuEST 光子计数干涉仪（PRX 2025），把'时空噪声'变成有灵敏度的实验指标", 13, False, DARK),
    ("· 核钟：Th-229 直接激光激发（148 nm）+ 无自旋宿主晶体，neV 窗口的探测器已就位", 13, False, DARK),
    ("同时，'统一场论'这个词在大众传播中仍然与反重力传说绑定。", 14, False, DARK),
    ("本论文的时点价值：把'什么命题还活着、用什么实验判决、花多少钱'第一次做成一张可复现的地图。", 14, True, ACCENT),
])
rows = [
    ["时间", "里程碑", "对判决矩阵的意义"],
    ["2024", "Panda 格点原子干涉 6.2 nm/s²", "第 13 行：屏蔽第五力排除"],
    ["2025", "GQuEST PRX / Aziz–Howl Nature 争论", "第 15 行通道 + 第 6/7 行方法论"],
    ["2026", "Th-229 无自旋宿主 / Céleri 反冲约束", "第 12 行窗口 + 叠加制备新边界"],
]
add_table(s, rows, Inches(7.4), Inches(1.3), Inches(5.4), Inches(2.6), font_size=11.5)
add_footer(s, "来源：实验路线图.md 第七十六节", 2)
add_notes(s, "强调时点性：三个时间尺度上的实验进展把'量子引力'从哲学问题变成了工程问题。判决矩阵的每一行都对应一个正在进行的实验。")

# ================= Slide 3: 三种统一（概念框架） =================
s = new_slide()
add_title(s, "概念框架：三种统一分类", "03")
add_text(s, Inches(0.55), Inches(1.25), Inches(12.2), Inches(1.0), [
    ("'统一场论'是三个不同主张的混用。它们证据不同、失败方式不同、判决实验不同。",
     16, True, DARK),
    ("这是全篇的组织原则，也是本文的第一个原创贡献。", 13, False, GRAY),
])
rows = [
    ["类型", "命题", "失败方式", "判决通道"],
    ["转换式\nTransformational", "A 场物理地转化为 B 场\n（Tesla / 张祥前路线）", "定量失败\n（初等算术）", "Schwinger 极限反推\n+ 第五力全景"],
    ["几何式\nGeometrical", "A、B 是高维几何的投影\n（Kaluza–Klein）", "窗口收窄\n（未被排除）", "轨道约束 + Th-229\n核钟 neV 窗口"],
    ["代数式\nAlgebraic", "A、B 共享更深代数\n（double copy）", "本体论沉默\n（数学成立）", "不可直接判决\n（归纳问题）"],
]
add_table(s, rows, Inches(0.55), Inches(2.3), Inches(12.2), Inches(3.6),
          col_widths=[Inches(1.8), Inches(4.6), Inches(2.8), Inches(3.0)], font_size=13)
add_text(s, Inches(0.55), Inches(6.15), Inches(12.2), Inches(0.7), [
    ("要点：三个断言类别，而非三个候选理论。混淆它们，是大众传播里百年混乱的根源。",
     14, True, ACCENT)])
add_footer(s, "来源：核心命题-形式化.md；投稿-概念论文.md", 3)
add_notes(s, "讲解三种分类时各举一个例子：转换式=特斯拉传说；几何式=Kaluza-Klein；代数式=double copy。强调'失败方式不同'这个判据比'对错'更有信息量。")

# ================= Slide 4: Reality Stack =================
s = new_slide()
add_title(s, "全局地图：Reality Stack 七层", "04")
add_pic(s, "fig02_reality_stack.png", Inches(2.3), Inches(1.15), h=Inches(5.2))
add_text(s, Inches(0.55), Inches(6.5), Inches(12.2), Inches(0.5), [
    ("人类在测量栈上远远领先，在控制栈上几乎空白——这正是'Gravity is the wrong engineering variable'的含义。",
     13, True, BLUE)], PP_ALIGN.CENTER)
add_footer(s, "图 2：Reality Stack（figures/02_reality_stack.py 生成）", 4)
add_notes(s, "七层从观测量到未知结构。重点讲 Layer 4 的'唯一入口'标注：在 GR+SM 内，一切工程操作对运动的影响都由它对 Tμν 的影响决定。")

# ================= Slide 5: 主题1 转换式统一 =================
s = new_slide()
add_title(s, "主题一：转换式统一已被定量证伪", "05")
add_text(s, Inches(0.55), Inches(1.3), Inches(6.9), Inches(5.4), [
    ("问题：变化的电磁场能产生引力吗？", 17, True, DARK),
    ("能——Tμν 里既有能量又有压强（因子 2），但量级如何？", 14, False, DARK),
    ("正面：1 m³、10 MV/m 的电容器产生 ~10⁻²⁵ m/s² 引力，比地球小 25 个数量级。", 14, False, DARK),
    ("反推：要产生地球级引力，需要场强达到 Schwinger 极限的 ~8 倍——该处真空已不稳定，你制造的是物质而非反重力。", 14, False, ACCENT),
    ("结论：转换式统一的实验室版本被两条算术关闭，无需任何新物理。", 14, True, DARK),
    ("", 6, False, DARK),
    ("反推公式：", 13, True, GRAY),
    ("E_needed = √(2ρ/ε₀)，ρ = g c²/(8π G L)", 13, False, BLUE),
])
add_pic(s, "fig04_fifth_force.png", Inches(7.6), Inches(1.35), w=Inches(5.3))
add_footer(s, "图：第五力约束全景；模拟 02（EM 场引力效应）", 5)
add_notes(s, "讲清楚'因子 2'：压强也进 Tμν，所以电磁场的引力源是其能量密度的两倍——这是论文第十二节'关键变量是压力'的数值演示。Schwinger 反推是全文最强的段落。")

# ================= Slide 6: 主题2 几何式统一 =================
s = new_slide()
add_title(s, "主题二：几何式统一的 neV 窗口", "06")
add_text(s, Inches(0.55), Inches(1.3), Inches(6.6), Inches(5.4), [
    ("Kaluza–Klein：额外维的引力子模式给出 Yukawa 修正", 16, True, DARK),
    ("V(r) = −(Gm₁m₂/r)(1 + α e^(−r/λ))", 14, False, BLUE),
    ("引力 AB 实验把它变成能级分裂：ΔE = m_sys · α · ΔΦ(r,λ)", 14, False, BLUE),
    ("文献宣称的 eV（核）/ meV（原子）级分裂需要 α ~ 0.4 / 10⁻³——被轨道第五力约束排除 4–8 个数量级。", 14, False, DARK),
    ("真正开放的窗口在 neV 量级，探测器硬件链已全部就位。", 14, True, GREEN),
])
rows = [
    ["核钟节点", "状态"],
    ["148.18 nm 直接激发", "已实现（Nat Commun 2024）"],
    ["CaF₂ 宿主寿命 641 s", "已测"],
    ["无自旋宿主 Th(SO₄)₂", "潜在稳定度 4.6×10⁻²³/√τ"],
    ["X 射线淬灭读出", "周期加速 ≥50 倍"],
]
add_table(s, rows, Inches(7.4), Inches(1.4), Inches(5.4), Inches(2.8), font_size=11.5)
add_text(s, Inches(7.4), Inches(4.5), Inches(5.4), Inches(1.5), [
    ("几何式统一不再是一个模糊希望：它是一个有确定仪器、确定能量窗口的测量计划。", 13, True, GREEN)])
add_footer(s, "来源：研究缺口补齐.md A 节；experiments/05、08", 6)
add_notes(s, "强调'窗口已知'：不是'额外维被证伪'，而是'可测窗口被压缩到 neV 且探测器就位'。这是对 Kaluza-Klein 最诚实的状态描述。")

# ================= Slide 7: 主题3 代数式统一 =================
s = new_slide()
add_title(s, "主题三：代数式统一——数学成立，本体论沉默", "07")
add_text(s, Inches(0.55), Inches(1.3), Inches(12.2), Inches(4.2), [
    ("double copy 的三个代数事实：", 16, True, DARK),
    ("· off-shell N=8 超引力 = (N=4 超杨-米尔斯)²（三次阶）——Bonezzi–Casale–Hohm 2025", 14, False, DARK),
    ("· 相干态规范背景上的振幅 double copy 到弯曲时空，metric 由规范背景构造——Ilderton–Lindved 2025", 14, False, DARK),
    ("· Ehlers 变换 = 电磁对偶的 double copy——Alawadhi 2023（Tesla 梦寐以求的'EM↔引力关系'，其现代形态是代数平方）", 14, False, DARK),
    ("但边界同样明确：Schwarzschild 的剩余规范代数在引力侧坍缩为有限等距群（Holton 2025）——共享代数结构不等于共享物理内容。", 14, False, ACCENT),
    ("", 6, False, DARK),
    ("诚实结论：代数式统一是真实的数学结构，其本体论地位不可由计算判决——这是对 double copy 能声称什么的最精确边界。", 14, True, DARK),
])
add_footer(s, "来源：时空之前是什么.md 第六十二节", 7)
add_notes(s, "三个事实逐个过一遍，重点落在 Ehlers=电磁对偶的平方：这是把 Tesla 的问题从'转换'翻译成'代数'的钥匙。Holton 的坍缩结果是防止过度解读的刹车。")

# ================= Slide 8: 争议 GIE =================
s = new_slide()
add_title(s, "争议：纠缠能否证明引力是量子的？", "08")
add_text(s, Inches(0.55), Inches(1.3), Inches(6.9), Inches(5.2), [
    ("2025 年 Nature 争论的时间线：", 15, True, DARK),
    ("① BMV 逻辑：两质量只经引力作用产生纠缠 ⇒ 中介非经典", 13, False, DARK),
    ("② Aziz–Howl：'经典引力 + 量子场物质'也能产生纠缠", 13, False, DARK),
    ("③ Marletto 等反驳：非相对论极限下相互作用超局域，无纠缠", 13, False, DARK),
    ("④ 2026 解剖：Aziz–Howl 效应是物质扇区串扰（不同物质种类时归零）", 13, False, DARK),
    ("方法论结论：GIE 是 LOCC 排除器，不是'一切经典'的排除器；", 14, True, ACCENT),
    ("需要一个理论歧视树，而非一锤定音实验。", 14, True, ACCENT),
    ("", 6, False, DARK),
    ("我们的数值实现：量子引力 N≈0.078（CHSH 违背 0.024），平均场/LOCC 严格为零；", 13, False, DARK),
    ("qubit 抽象忠实性误差 ~e^(−2α²)（α=2 时 ~10⁻⁴）。", 13, False, DARK),
])
add_pic(s, "fig05_gie_discrimination.png", Inches(7.7), Inches(1.35), w=Inches(5.2))
add_footer(s, "图：GIE 三模型歧视（模拟 01/06）；来源：第六十四节", 8)
add_notes(s, "这是'箭头纪律'的活例子：观测到 O 不能证明理论 A，如果理论 B 也能产生 O。争议本身比结论更有价值——它教会我们什么才算判决。")

# ================= Slide 9: 判决矩阵 =================
s = new_slide()
add_title(s, "判决矩阵：19 个箭头的当前状态", "09")
add_pic(s, "fig01_decision_matrix.png", Inches(3.3), Inches(1.1), h=Inches(5.3))
add_text(s, Inches(0.55), Inches(6.55), Inches(12.2), Inches(0.5), [
    ("统计：4 不成立 + 2 成立（电与磁、电弱）+ 1 分类陈述 + 4 个 A 级事实 + 9 行 B/C 级（均有判决对象与时间尺度）",
     12, True, DARK)], PP_ALIGN.CENTER)
add_footer(s, "图 1：判决矩阵全图（19 行）；来源：实验路线图.md 第八十三节", 9)
add_notes(s, "整页给图。强调方法论：每行必须携带判决通道+判决对象+时间尺度，16/19 行有可复现脚本支撑。这是论文的最终交付物。")

# ================= Slide 10: 中心定理 =================
s = new_slide()
add_title(s, "中心定理：工程控制通道唯一性", "10")
add_text(s, Inches(0.55), Inches(1.35), Inches(12.2), Inches(1.1), [
    ("在 GR + SM 内，若工程操作满足 (A1) 源不变、(A2) 度规不被直接修改、(A3) 无新通道，则它不改变任何测地线运动。",
     16, True, DARK),
])
boxes = [
    ("违反 A1", "改变 Tμν", "能量/压强/流工程——GR 允许的全部'惯性工程'，即普通能量工程", GREEN),
    ("违反 A2", "直接写度规", "度规工程——需非常规物质或修改引力（warp-drive no-go 文献是其特殊情形）", BLUE),
    ("违反 A3", "引入新通道", "第五力 / LIV / GUP——判决矩阵第 13、19 行，已被大量排除", PURPLE),
]
x0 = 0.55
for i, (head, mid, body, color) in enumerate(boxes):
    x = Inches(x0 + i * 4.15)
    box = s.shapes.add_shape(1, x, Inches(2.7), Inches(3.9), Inches(2.6))
    box.fill.solid(); box.fill.fore_color.rgb = LIGHT
    box.line.color.rgb = color; box.line.width = Pt(1.6)
    add_text(s, x + Inches(0.15), Inches(2.85), Inches(3.6), Inches(0.5), [(head, 15, True, color)])
    add_text(s, x + Inches(0.15), Inches(3.35), Inches(3.6), Inches(0.5), [(mid, 14, True, DARK)])
    add_text(s, x + Inches(0.15), Inches(3.95), Inches(3.6), Inches(1.3), [(body, 11.5, False, DARK)])
add_text(s, Inches(0.55), Inches(5.7), Inches(12.2), Inches(1.1), [
    ("推论 2（杠杆效力）：质量分解的每一成分作为引力工程杠杆的效力 = ∂Tμν/∂该成分——改 Higgs 只动 ~9% 的质子质量项，改 QCD 动 ~91%，改'真空'在定义层失败。",
     14, True, ACCENT),
    ("推论 3：超材料负有效质量是色散关系工程，不改 Tμν——与引力无关。", 14, True, ACCENT),
])
add_footer(s, "来源：核心命题-形式化.md", 10)
add_notes(s, "定理本身近乎同义反复——价值在假设的完备性：(A1)(A2)(A3) 穷尽 GR 内所有控制通道。推论 2 是最接近原创的单个命题。诚实声明：这是分类性陈述，不是新现象预言。")

# ================= Slide 11: 可复现性 =================
s = new_slide()
add_title(s, "可复现性：每个箭头都有数字", "11")
rows = [
    ["箭头", "关键数字", "含义"],
    ["GIE 纠缠相位", "φ = 0.217 rad（基准）vs 10⁻⁶~10⁻⁵（芯片参数）", "3σ 认证从 200 次 → 10⁹–10²¹ 次"],
    ["geontropic 涨落", "强档 0.03 s vs 弱档 10³³ s", "可证伪窗口只在强档"],
    ["frame dragging", "距读出能力 ~7 个数量级", "10 年+，缺口主要在源"],
    ["KK 分裂", "eV 窗口排除 4–8 个数量级", "neV 窗口由核钟接管"],
    ["装置反冲（Céleri）", "1 kg@1K 装置封顶 10⁹ u 叠加于 ~100 μm", "MAQRO 量级的真实约束"],
]
add_table(s, rows, Inches(0.55), Inches(1.4), Inches(12.2), Inches(3.9),
          col_widths=[Inches(2.6), Inches(5.4), Inches(4.2)], font_size=13)
add_text(s, Inches(0.55), Inches(5.6), Inches(12.2), Inches(1.0), [
    ("16 份脚本（9 核算 + 7 模拟）全部通过；两张关键图（判决矩阵、第五力全景）直接由脚本生成。",
     14, True, DARK),
    ("'核算之前和之后，一条箭头是两个不同的对象'——论文的全部方法论就是把箭头变成可称量的对象。", 14, True, BLUE),
])
add_footer(s, "来源：experiments/ 与 simulations/ 全部脚本", 11)
add_notes(s, "逐行读表，强调每个数字对应一个判决矩阵行。Céleri 的'装置反冲'是 2026 年新约束，说明叠加制备的极限不只是环境退相干。")

# ================= Slide 12: RT 极小面 =================
s = new_slide()
add_title(s, "纠缠 → 几何：RT 极小面的玩具实现", "12")
add_pic(s, "fig03_rt_mincut.png", Inches(1.6), Inches(1.3), w=Inches(7.2))
add_text(s, Inches(9.1), Inches(1.5), Inches(3.7), Inches(4.6), [
    ("等距树张量网络中：", 14, True, DARK),
    ("S₂(A) ≈ mincut(A)·lnχ", 14, False, BLUE),
    ("· 半树与相邻两叶同属 mincut=1 类——熵由极小面决定，而非 |A|", 12, False, DARK),
    ("· 远离两叶给出 ~2 倍熵——'断开的极小面'（玩具层的相变骨架）", 12, False, DARK),
    ("· 手写树与 quimb 原生两种独立实现互验", 12, False, DARK),
    ("· O(1/χ²) 有限键修正 = 大 N 极限的玩具对应", 12, False, DARK),
    ("保真度声明：机制演示，不是'真实宇宙由张量网络构成'的证据（判决矩阵第 18 行，C 级）。", 12, True, ACCENT),
])
add_footer(s, "图 3：RT 极小面互验（模拟 05/07）", 12)
add_notes(s, "这是第 18 行'纠缠→几何'的玩具骨架。强调边界：全息对偶内部是严格的，对真实宇宙仍是 C 级假说。")

# ================= Slide 13: 未来方向 =================
s = new_slide()
add_title(s, "未来方向：近 / 中 / 远期行动清单", "13")
cols = [
    ("近期 3–5 年", GREEN, [
        "GQuEST 首轮数据（第 15 行）",
        "reservoir-engineered GIE（Q 因子降一个量级）",
        "CSS WEP 系统误差账本",
        "钍核钟 10⁻¹⁸ 稳定度",
        "Céleri 反冲约束实验验证",
    ]),
    ("中期 5–10 年", BLUE, [
        "QGEM 首轮（纳米金刚石 SGI）",
        "引力 AB 的区分几何设计",
        "LMT 10⁻¹⁷ WEP（STE-QUEST 规格）",
        "MAQRO 级空间任务（坍缩模型判决）",
    ]),
    ("远期 10 年+", PURPLE, [
        "frame-dragging GIE（缺口 ~7 量级）",
        "bolometric 引力子噪声",
        "量子因果结构（Hardy QEP）",
    ]),
]
for i, (head, color, items) in enumerate(cols):
    x = Inches(0.55 + i * 4.2)
    box = s.shapes.add_shape(1, x, Inches(1.3), Inches(3.9), Inches(5.1))
    box.fill.solid(); box.fill.fore_color.rgb = LIGHT
    box.line.color.rgb = color; box.line.width = Pt(1.4)
    add_text(s, x + Inches(0.12), Inches(1.4), Inches(3.7), Inches(0.5), [(head, 15, True, color)])
    lines = [(("· " + it), 11.5, False, DARK) for it in items]
    add_text(s, x + Inches(0.12), Inches(1.95), Inches(3.7), Inches(4.2), lines)
add_footer(s, "来源：物理现实的工程学边界.md 第八十八节（12 条行动清单）", 13)
add_notes(s, "每一条都有判决对象和时间尺度——这就是'路线图'与'愿望清单'的区别。重点提 QGEM 首轮与核钟，它们是最可能改变矩阵的两条。")

# ================= Slide 14: 总结 =================
s = new_slide()
dark_bg(s)
add_text(s, Inches(1.2), Inches(1.1), Inches(11), Inches(0.7),
         [("总结：核心命题三层", 28, True, WHITE)], PP_ALIGN.CENTER)
layers = [
    ("历史层", "Tesla 与张祥前式的'转换式统一'没有实验依据，但他们的提问——世界背后是否有统一动力结构——是现代基础物理的真实问题。"),
    ("物理层", "Gravity 是错误工程变量；Inertia 也是；唯一入口是 Tμν，而 Tμν 的正确工程名字叫'能量工程'。"),
    ("前沿层", "Tμν 本身可能是派生的。真正的源变量候选：量子态、运动学代数、关系结构——每个都有已排期的实验。"),
]
for i, (head, body) in enumerate(layers):
    y = Inches(1.95 + i * 1.45)
    add_text(s, Inches(1.5), y, Inches(1.6), Inches(0.6), [(head, 17, True, ACCENT)])
    add_text(s, Inches(3.2), y, Inches(8.7), Inches(1.2), [(body, 14, False, WHITE)])
add_text(s, Inches(1.2), Inches(6.35), Inches(11), Inches(0.6), [
    ("反重力如果存在，不会是一台抵抗引力的机器——而地图的第一步已经完成：19 行，16 份脚本，全部可复现。",
     15, True, RGBColor(0xF5, 0xC6, 0x42))], PP_ALIGN.CENTER)
add_footer(s, "", 14)
add_notes(s, "收束：三层命题是全文的压缩。最后一句点出'地图是交付物'——本论文不预言新现象，它把模糊野心变成可检验命题。")

# ================= Slide 15: 参考文献与代码 =================
s = new_slide()
add_title(s, "参考文献与代码", "15")
add_text(s, Inches(0.55), Inches(1.35), Inches(12.2), Inches(4.6), [
    ("全部可复现资产：github.com/everest-an/Antigravity", 16, True, DARK),
    ("· 论文全稿：90 节 + 严谨化附件 A–G（中英双语，PDF 版在 build/ 目录）", 13, False, DARK),
    ("· 核算 9 份 + 模拟 7 份（run_all.py 一键全跑，0 失败）", 13, False, DARK),
    ("· 投稿图 5 幅（PNG + SVG，脚本生成）", 13, False, DARK),
    ("· 参考文献总表 ~110 条（含预印本/同行评议状态标注）", 13, False, DARK),
    ("", 6, False, DARK),
    ("核心文献：Bose et al. PRL 2017 · Fedida–Kent PRD 2025 · Salzger–Vilasini arXiv:2605.08351 · GQuEST PRX 2025 · Derevianko et al. Colloquium 2026 · Bobrick–Martire CQG 2021 · Barzegar et al. arXiv:2602.16495", 11.5, False, GRAY),
    ("", 6, False, DARK),
    ("致谢：本汇报按 nature-skills（Yuan1z0825）的 nature-paper2ppt 工作流生成。", 11, False, GRAY),
])
add_footer(s, "", 15)
add_notes(s, "收尾页。留时间回答提问：最常见的问题是'这是不是反重力论文'——回答：不是，这是关于'为什么反重力说法全部失败以及真正开放的问题在哪里'的地图。")

OUT.mkdir(exist_ok=True)
out_path = OUT / "从特斯拉到时空工程-汇报.pptx"
prs.save(str(out_path))
print("saved:", out_path, "slides:", len(prs.slides._sldIdLst))
